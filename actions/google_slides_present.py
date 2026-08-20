"""
google_slides_present.py — Present Google Slides in Chrome with JARVIS narration.

Opens the /present URL in the user's real Chrome profile (Google login intact),
captures each slide from the screen, and advances with keyboard control.
"""
from __future__ import annotations

import hashlib
import platform
import re
import subprocess
import threading
import time
from pathlib import Path
from typing import Callable, Optional
from urllib.parse import urlparse

_OS = platform.system()
_CONFIG_PATH = Path(__file__).resolve().parent.parent / "config" / "api_keys.json"


class PresentationState:
    """Shared state so stop/interrupt can cancel an in-flight presentation."""

    def __init__(self) -> None:
        self._stop = threading.Event()
        self._lock = threading.Lock()
        self.running = False
        self.slide_num = 0
        self.present_url = ""
        self.deck_id = ""
        self.slide_index: dict[str, list[dict]] = {}  # deck_id -> [{num, text}, ...]

    def request_stop(self) -> None:
        self._stop.set()

    def clear_stop(self) -> None:
        self._stop.clear()

    def should_stop(self) -> bool:
        return self._stop.is_set()

    def mark_running(self, url: str) -> None:
        with self._lock:
            self.running = True
            self.slide_num = 0
            self.present_url = url
            self.deck_id = extract_deck_id(url)
            self._stop.clear()

    def mark_stopped(self) -> None:
        with self._lock:
            self.running = False
            self.present_url = ""

    def set_slide(self, n: int) -> None:
        with self._lock:
            self.slide_num = n

    def get_index(self, deck_id: str) -> list[dict]:
        with self._lock:
            return list(self.slide_index.get(deck_id, []))

    def save_index(self, deck_id: str, index: list[dict]) -> None:
        with self._lock:
            self.slide_index[deck_id] = index


_STATE = PresentationState()


def get_state() -> PresentationState:
    return _STATE


def stop_presentation() -> str:
    _STATE.request_stop()
    return "Presentation stop requested."


def parse_present_url(url: str) -> str:
    """Normalize any Google Slides URL to fullscreen present mode."""
    url = (url or "").strip()
    if not url:
        raise ValueError("No Google Slides URL provided.")

    deck_id = extract_deck_id(url)
    if not deck_id:
        raise ValueError(
            "Invalid Google Slides URL. Expected "
            "https://docs.google.com/presentation/d/ID/edit"
        )
    return f"https://docs.google.com/presentation/d/{deck_id}/present"


def extract_deck_id(url: str) -> str:
    match = re.search(r"/presentation/d/([a-zA-Z0-9_-]+)", url or "")
    return match.group(1) if match else ""


def _get_api_key() -> str:
    try:
        import json
        return json.loads(_CONFIG_PATH.read_text(encoding="utf-8")).get("gemini_api_key", "")
    except Exception:
        return ""


def build_slide_index_from_pptx(path: Path) -> list[dict]:
    """Extract slide number + text from a local .pptx export."""
    from pptx import Presentation

    prs = Presentation(path)
    index: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        index.append({"num": i, "text": "\n".join(parts)})
    return index


def find_slide_in_index(topic: str, index: list[dict]) -> Optional[int]:
    """Find first slide whose text mentions topic (e.g. SLO)."""
    needle = topic.lower().strip()
    if not needle:
        return None

    for entry in index:
        hay = entry.get("text", "").lower()
        if needle in hay:
            return int(entry["num"])

    # Token match — e.g. topic "service level" in longer bullet text
    tokens = [t for t in re.split(r"\W+", needle) if len(t) > 2]
    for entry in index:
        hay = entry.get("text", "").lower()
        if tokens and all(tok in hay for tok in tokens):
            return int(entry["num"])

    return None


def _gemini_slide_matches_topic(img_bytes: bytes, mime: str, topic: str) -> bool:
    """Vision check: does the visible slide mention this topic?"""
    import base64

    key = _get_api_key()
    if not key:
        return False
    try:
        from google import genai
        client = genai.Client(api_key=key)
        b64 = base64.b64encode(img_bytes).decode("ascii")
        prompt = (
            f'Does this presentation slide discuss, define, or prominently mention '
            f'"{topic}" (or its acronym)? Reply with ONLY the word yes or no.'
        )
        models = ("gemini-2.0-flash", "gemini-2.5-flash-lite", "gemini-3.6-flash")
        last_err = None
        for model in models:
            try:
                resp = client.models.generate_content(
                    model=model,
                    contents=[
                        {"inline_data": {"mime_type": mime, "data": b64}},
                        prompt,
                    ],
                )
                answer = (resp.text or "").strip().lower()
                return answer.startswith("yes")
            except Exception as e:
                last_err = e
                if "404" in str(e) or "NOT_FOUND" in str(e):
                    continue
                raise
        if last_err:
            raise last_err
    except Exception as e:
        print(f"[SlidesPresent] Vision topic match failed: {e}")
    return False


def go_to_slide_number(target: int) -> None:
    """Jump to slide N from the beginning (Home + Right × N-1)."""
    target = max(1, int(target))
    focus_chrome()
    press_key("home")
    time.sleep(0.55)
    _STATE.set_slide(1)
    for _ in range(target - 1):
        press_key("right")
        time.sleep(0.35)
    _STATE.set_slide(target)


def find_slide_by_topic_vision(
    topic: str,
    capture_fn: Callable[[], tuple[bytes, str]],
    *,
    max_slides: int = 40,
) -> Optional[int]:
    """Scan from slide 1 until a slide visually matches the topic."""
    focus_chrome()
    press_key("home")
    time.sleep(0.8)
    _STATE.set_slide(1)

    for n in range(1, max_slides + 1):
        if _STATE.should_stop():
            break
        img_bytes, mime = capture_fn()
        if _gemini_slide_matches_topic(img_bytes, mime, topic):
            _STATE.set_slide(n)
            return n
        if n < max_slides:
            press_key("right")
            time.sleep(0.65)
            _STATE.set_slide(n + 1)
    return None


def go_to_topic(
    topic: str,
    *,
    deck_id: str = "",
    pptx_path: str = "",
    url: str = "",
    capture_fn: Optional[Callable[[], tuple[bytes, str]]] = None,
    max_slides: int = 40,
) -> str:
    """Find a slide about topic and jump to it."""
    topic = topic.strip()
    if not topic:
        return "Provide a topic to search for (e.g. SLO)."

    did = deck_id or extract_deck_id(url) or _STATE.deck_id
    index = _STATE.get_index(did) if did else []

    if not index and pptx_path:
        pptx = Path(pptx_path).expanduser()
        if pptx.exists():
            index = build_slide_index_from_pptx(pptx)
            if did:
                _STATE.save_index(did, index)

    if index:
        num = find_slide_in_index(topic, index)
        if num:
            go_to_slide_number(num)
            return f"Found '{topic}' on slide {num} (from deck index)."
        return f"No slide in the indexed deck mentions '{topic}'."

    if not capture_fn:
        return (
            f"No slide index for '{topic}'. Keep the slideshow open in Chrome, or export "
            "the deck as .pptx and pass pptx_path for faster topic search."
        )

    num = find_slide_by_topic_vision(topic, capture_fn, max_slides=max_slides)
    if num:
        return f"Found '{topic}' on slide {num} (visual scan)."
    return f"Could not find a slide about '{topic}' in the first {max_slides} slides."


def open_in_chrome(url: str) -> str:
    """Open URL in Google Chrome using the user's logged-in profile."""
    if _OS == "Darwin":
        subprocess.run(
            ["open", "-a", "Google Chrome", url],
            check=False,
            timeout=20,
        )
        return f"Opened in Google Chrome: {url}"
    if _OS == "Windows":
        chrome_paths = [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        ]
        for exe in chrome_paths:
            try:
                subprocess.Popen([exe, url], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                return f"Opened in Chrome: {url}"
            except Exception:
                continue
        subprocess.Popen(f'start chrome "{url}"', shell=True)
        return f"Opened in Chrome: {url}"
    # Linux
    for cmd in ("google-chrome", "google-chrome-stable", "chromium-browser", "chromium"):
        try:
            subprocess.Popen([cmd, url], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return f"Opened in {cmd}: {url}"
        except Exception:
            continue
    raise RuntimeError("Could not find Chrome on this system.")


def focus_chrome() -> None:
    if _OS == "Darwin":
        subprocess.run(
            ["osascript", "-e", 'tell application "Google Chrome" to activate'],
            check=False,
            timeout=5,
        )
    elif _OS == "Windows":
        try:
            import pygetwindow as gw
            for w in gw.getAllWindows():
                if "chrome" in (w.title or "").lower():
                    w.activate()
                    return
        except Exception:
            pass
    time.sleep(0.3)


def press_key(key: str) -> None:
    try:
        import pyautogui
        focus_chrome()
        time.sleep(0.2)
        pyautogui.press(key)
    except Exception as e:
        raise RuntimeError(f"Keyboard control failed: {e}") from e


def advance_slide() -> None:
    press_key("right")
    with _STATE._lock:
        _STATE.slide_num = max(1, _STATE.slide_num + 1)


def previous_slide() -> None:
    press_key("left")
    with _STATE._lock:
        _STATE.slide_num = max(1, _STATE.slide_num - 1)


def exit_present_mode() -> None:
    """Escape exits Google Slides present mode."""
    try:
        focus_chrome()
        press_key("escape")
    except Exception:
        pass


def _image_hash(img_bytes: bytes) -> str:
    return hashlib.md5(img_bytes).hexdigest()


def run_presentation_loop(
    url: str,
    *,
    max_slides: int = 25,
    load_wait: float = 8.0,
    slide_wait: float = 1.2,
    capture_fn: Callable[[], tuple[bytes, str]],
    present_slide_fn: Callable[[bytes, str, int], None],
    wait_after_speech_fn: Callable[[], None],
    on_status: Optional[Callable[[str], None]] = None,
    on_slide: Optional[Callable[[int, str], None]] = None,
) -> str:
    """
    Synchronous presentation loop — call from asyncio executor or background thread.

    present_slide_fn(img, mime, slide_num) should trigger JARVIS to speak about the slide.
    wait_after_speech_fn() should block until narration audio finishes.
    """
    if _STATE.running:
        _STATE.request_stop()
        time.sleep(0.5)

    try:
        present_url = parse_present_url(url)
    except ValueError as e:
        return str(e)

    _STATE.mark_running(present_url)
    _log = on_status or (lambda m: print(f"[SlidesPresent] {m}"))

    try:
        _log("Opening Google Slides present mode in Chrome…")
        open_in_chrome(present_url)
        _log(f"Waiting {load_wait:.0f}s for slideshow to load — keep Chrome in front.")
        time.sleep(load_wait)
        focus_chrome()

        prev_hash = ""
        same_count = 0

        for slide_num in range(1, max_slides + 1):
            if _STATE.should_stop():
                _log("Presentation stopped.")
                break

            _STATE.set_slide(slide_num)
            _log(f"Slide {slide_num} — capturing screen…")

            try:
                img_bytes, mime = capture_fn()
            except Exception as e:
                return f"Screen capture failed on slide {slide_num}: {e}"

            img_hash = _image_hash(img_bytes)
            if img_hash == prev_hash:
                same_count += 1
                if same_count >= 2:
                    _log("End of deck detected (duplicate slide).")
                    break
            else:
                same_count = 0
            prev_hash = img_hash

            if on_slide:
                on_slide(slide_num, f"Presenting slide {slide_num}")

            try:
                present_slide_fn(img_bytes, mime, slide_num)
                wait_after_speech_fn()
            except Exception as e:
                _log(f"Narration error on slide {slide_num}: {e}")

            if _STATE.should_stop():
                break

            if slide_num < max_slides:
                _log(f"Advancing to slide {slide_num + 1}…")
                advance_slide()
                time.sleep(slide_wait)

        return f"Presentation finished after { _STATE.slide_num } slide(s)."

    except Exception as e:
        return f"Presentation failed: {e}"

    finally:
        exit_present_mode()
        _STATE.mark_stopped()


def google_slides_present(
    parameters: dict | None = None,
    player=None,
    *,
    capture_fn=None,
    present_slide_fn=None,
    wait_after_speech_fn=None,
) -> str:
    """
    Tool entry point. Actions: start | stop | next

    For 'start', the caller (main.py) should run the loop asynchronously with callbacks.
    """
    params = parameters or {}
    action = (params.get("action") or "start").lower().strip()

    if action == "stop":
        return stop_presentation()

    if action in ("next", "forward", "advance"):
        try:
            advance_slide()
            return f"Advanced to slide {_STATE.slide_num}."
        except Exception as e:
            return str(e)

    if action in ("previous", "back", "prior", "last"):
        try:
            previous_slide()
            return f"Went back to slide {_STATE.slide_num}."
        except Exception as e:
            return str(e)

    if action == "status":
        if _STATE.running:
            return f"Presenting slide {_STATE.slide_num}: {_STATE.present_url}"
        if _STATE.slide_num:
            return f"On slide {_STATE.slide_num}."
        return "No presentation is running."

    if action in ("goto", "go_to", "jump", "find"):
        slide_number = params.get("slide_number")
        topic = (
            params.get("topic")
            or params.get("query")
            or params.get("subject")
            or params.get("text")
            or ""
        ).strip()
        pptx_path = (params.get("pptx_path") or "").strip()
        url = (params.get("url") or _STATE.present_url or "").strip()

        if slide_number is not None and str(slide_number).strip() != "":
            try:
                go_to_slide_number(int(slide_number))
                return f"Jumped to slide {int(slide_number)}."
            except Exception as e:
                return str(e)

        if topic:
            try:
                return go_to_topic(
                    topic,
                    deck_id=extract_deck_id(url),
                    pptx_path=pptx_path,
                    url=url,
                    capture_fn=capture_fn,
                    max_slides=min(int(params.get("max_slides") or 40), 50),
                )
            except Exception as e:
                return f"Go-to topic failed: {e}"

        return "Provide slide_number (e.g. 5) or topic/query (e.g. SLO)."

    url = (params.get("url") or "").strip()
    if action != "start" and not url:
        # goto/start may omit url if slideshow already open
        if action == "start":
            return "Provide a Google Slides URL (action=start)."
        url = _STATE.present_url

    if action != "start":
        return f"Unknown action: {action}"

    if not url:
        return "Provide a Google Slides URL (action=start)."

    if not all([capture_fn, present_slide_fn, wait_after_speech_fn]):
        return (
            "Presentation queued. "
            f"URL: {parse_present_url(url)} — "
            "the host must run the presentation loop with speech callbacks."
        )

    max_slides = min(int(params.get("max_slides") or 25), 50)
    load_wait = float(params.get("load_wait") or 8.0)

    def _status(msg: str) -> None:
        print(f"[SlidesPresent] {msg}")
        if player and hasattr(player, "write_log"):
            player.write_log(f"SLIDES: {msg}")

    def _on_slide(n: int, label: str) -> None:
        if player and hasattr(player, "show_content"):
            player.show_content(label, f"Google Slides — slide {n}")

    return run_presentation_loop(
        url,
        max_slides=max_slides,
        load_wait=load_wait,
        capture_fn=capture_fn,
        present_slide_fn=present_slide_fn,
        wait_after_speech_fn=wait_after_speech_fn,
        on_status=_status,
        on_slide=_on_slide,
    )
