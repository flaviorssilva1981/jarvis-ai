"""
presentation_create.py — Professional local .pptx decks (Gemini + python-pptx).

Creates styled PowerPoint files in ~/Documents/JARVIS Presentations/.
User uploads manually to Google Drive / Slides when ready.
"""
from __future__ import annotations

import io
import json
import platform
import re
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path

_BASE = Path(__file__).resolve().parent.parent
_CONFIG = _BASE / "config" / "api_keys.json"
_ASSETS = _BASE / "assets"
_LOGO_PATH = _ASSETS / "jarvis-youtube-thumbnail.png"
_OUTPUT_DIR = Path.home() / "Documents" / "JARVIS Presentations"
_GEMINI_MODELS = ("gemini-3.6-flash", "gemini-2.5-flash-lite", "gemini-2.0-flash")
_IMAGE_MODELS = ("imagen-3.0-generate-002", "imagen-3.0-fast-generate-001")
_BUILD_WITH_IMAGES = False  # Imagen often unavailable; PIL diagrams are the primary visuals

# Corporate JARVIS theme
_THEME = {
    "bg_dark": (0x0B, 0x1F, 0x3A),
    "bg_light": (0xF4, 0xF7, 0xFA),
    "accent": (0x00, 0xB4, 0xD8),
    "accent2": (0x00, 0x7A, 0x99),
    "white": (0xFF, 0xFF, 0xFF),
    "text_dark": (0x1A, 0x1A, 0x2E),
    "text_muted": (0x5A, 0x6A, 0x7A),
}


def _rgb(name: tuple[int, int, int]):
    from pptx.dml.color import RGBColor
    return RGBColor(*name)


def _get_api_key() -> str:
    return json.loads(_CONFIG.read_text(encoding="utf-8"))["gemini_api_key"]


def _gemini_generate(prompt: str) -> str:
    from google import genai

    client = genai.Client(api_key=_get_api_key())
    last_err = None
    for model in _GEMINI_MODELS:
        try:
            resp = client.models.generate_content(model=model, contents=prompt)
            return (resp.text or "").strip()
        except Exception as e:
            last_err = e
            err = str(e)
            if "404" in err or "NOT_FOUND" in err or "no longer available" in err.lower():
                print(f"[PresentationCreate] Model {model} unavailable — trying next…")
                continue
            raise
    raise last_err or RuntimeError("Gemini request failed.")


def _strip_json_fences(text: str) -> str:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text)
        text = re.sub(r"\s*```$", "", text)
    return text.strip()


def _safe_filename(title: str) -> str:
    name = re.sub(r"[^\w\s\-]", "", title).strip().replace(" ", "_")
    return (name[:60] or "presentation") + ".pptx"


def _generate_outline(
    title: str,
    topic: str,
    num_slides: int,
    audience: str = "",
) -> dict:
    audience_line = f"Audience: {audience}." if audience else ""
    prompt = f"""You are an expert presentation designer. Create a professional deck outline as JSON only.

Presentation title: {title}
Topic / brief: {topic}
Number of content slides (excluding title and closing): {num_slides}
{audience_line}

Return ONLY valid JSON:
{{
  "title": "Main title for title slide",
  "subtitle": "Compelling tagline",
  "theme_keywords": "technology, AI, professional",
  "slides": [
    {{
      "title": "Slide heading",
      "layout": "content",
      "bullets": ["Point one", "Point two"],
      "visual_hint": "Short description for a supporting illustration (icon/diagram style)"
    }},
    {{
      "title": "Architecture Overview",
      "layout": "diagram",
      "bullets": ["Brief caption"],
      "diagram_nodes": ["Voice Input", "Gemini AI", "Actions", "User"],
      "visual_hint": "System architecture flow"
    }},
    {{
      "title": "Key Takeaways",
      "layout": "closing",
      "bullets": ["Takeaway one", "Takeaway two"]
    }}
  ]
}}

Rules:
- Exactly {num_slides} items in "slides"
- layout: content | diagram | section | closing
- Include at least one "diagram" slide when topic is technical
- Last slide should be layout "closing"
- Bullets: concise, executive-ready (not paragraphs)
- visual_hint on every slide — describe a professional illustration
- diagram_nodes: 3-6 labels for diagram layout only
- No markdown, no code fences
"""
    raw = _gemini_generate(prompt)
    data = json.loads(_strip_json_fences(raw))
    if not isinstance(data.get("slides"), list) or not data["slides"]:
        raise ValueError("Gemini returned an invalid slide outline.")
    return _enhance_outline(data)


def _enhance_outline(outline: dict) -> dict:
    """Ensure technical decks get diagram/stack slides — not plain text placeholders."""
    for item in outline.get("slides", []):
        title = str(item.get("title") or "").lower()
        hint = str(item.get("visual_hint") or "").lower()
        layout = str(item.get("layout") or "content").lower()

        if any(k in title for k in ("architecture", "pipeline", "system flow", "how it works")):
            item["layout"] = "diagram"
            if not item.get("diagram_nodes"):
                item["diagram_nodes"] = _layers_from_bullets(item.get("bullets") or []) or [
                    "Voice / UI", "JARVIS Core", "Tools & APIs", "Output"
                ]
        elif any(k in title for k in ("stack", "technology", "components", "layers")):
            item["layout"] = "stack"
            item["stack_layers"] = _layers_from_bullets(item.get("bullets") or [])
        elif any(k in hint for k in ("pyramid", "stack", "layers")):
            item["layout"] = "stack"
            item["stack_layers"] = _layers_from_bullets(item.get("bullets") or [])

        if layout == "closing":
            item["layout"] = "closing"
    return outline


def _layers_from_bullets(bullets: list) -> list[str]:
    layers: list[str] = []
    for b in bullets:
        text = str(b).strip()
        if ":" in text:
            layers.append(text.split(":", 1)[0].strip())
        elif text:
            layers.append(text[:40])
    return layers[:6]


def _try_generate_slide_image(visual_hint: str, slide_title: str) -> Path | None:
    """Generate a slide illustration with Imagen (optional — skips on failure)."""
    if not _BUILD_WITH_IMAGES or not visual_hint:
        return None
    try:
        from google import genai
        from google.genai import types as gtypes

        client = genai.Client(api_key=_get_api_key())
        prompt = (
            f"Professional presentation illustration, flat modern corporate style, "
            f"clean vector look, dark blue and cyan palette, no text, no watermark. "
            f"Topic: {slide_title}. {visual_hint}"
        )
        for model in _IMAGE_MODELS:
            try:
                resp = client.models.generate_images(
                    model=model,
                    prompt=prompt,
                    config=gtypes.GenerateImagesConfig(number_of_images=1),
                )
                if resp.generated_images:
                    img_bytes = resp.generated_images[0].image.image_bytes
                    tmp = Path(tempfile.gettempdir()) / f"jarvis_slide_{abs(hash(visual_hint)) % 99999}.png"
                    tmp.write_bytes(img_bytes)
                    return tmp
            except Exception as e:
                print(f"[PresentationCreate] Image model {model}: {e}")
                continue
    except Exception as e:
        print(f"[PresentationCreate] Image generation skipped: {e}")
    return None


def _add_rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=None):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def _embed_png(slide, png_bytes: bytes, left: float, top: float, width: float | None = None, height: float | None = None):
    """Embed a Pillow-rendered PNG on a slide."""
    from pptx.util import Inches

    stream = io.BytesIO(png_bytes)
    if width is not None:
        slide.shapes.add_picture(stream, Inches(left), Inches(top), width=Inches(width))
    elif height is not None:
        slide.shapes.add_picture(stream, Inches(left), Inches(top), height=Inches(height))
    else:
        slide.shapes.add_picture(stream, Inches(left), Inches(top))


def _add_visual_png(slide, item: dict, left: float, top: float, width: float, height: float | None = None) -> None:
    """Render a professional diagram with Pillow and embed it."""
    from actions.presentation_graphics import pick_visual_png

    try:
        png = pick_visual_png(item)
        _embed_png(slide, png, left, top, width=width, height=height)
    except Exception as e:
        print(f"[PresentationCreate] PIL visual failed, using shapes: {e}")
        _draw_visual_panel(slide, item, left, top, width, height or 5.0)


def _add_logo(slide, left=11.8, top=0.25, height=0.55):
    from pptx.util import Inches

    logo = _LOGO_PATH if _LOGO_PATH.exists() else None
    if not logo:
        return
    try:
        slide.shapes.add_picture(str(logo), Inches(left), Inches(top), height=Inches(height))
    except Exception as e:
        print(f"[PresentationCreate] Logo insert failed: {e}")


def _add_decorative_shapes(slide, variant: str = "title"):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    accent = _rgb(_THEME["accent"])
    accent2 = _rgb(_THEME["accent2"])
    if variant == "title":
        _add_rect(slide, 0, 0, 13.333, 7.5, _rgb(_THEME["bg_dark"]))
        _add_rect(slide, 0, 6.85, 13.333, 0.08, accent)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(5.2), Inches(2.8), Inches(2.8))
        circ.fill.solid()
        circ.fill.fore_color.rgb = accent2
        circ.line.fill.background()
        circ2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.2), Inches(1.0), Inches(1.2), Inches(1.2))
        circ2.fill.solid()
        circ2.fill.fore_color.rgb = accent
        circ2.line.fill.background()
    elif variant == "content":
        _add_rect(slide, 0, 0, 13.333, 7.5, _rgb(_THEME["bg_light"]))
        _add_rect(slide, 0, 0, 0.18, 7.5, accent)
        _add_rect(slide, 0, 0, 13.333, 1.05, _rgb(_THEME["bg_dark"]))
        _add_rect(slide, 0, 1.05, 13.333, 0.06, accent)


def _add_title_slide(prs, outline: dict) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_decorative_shapes(slide, "title")
    _add_logo(slide)

    white = _rgb(_THEME["white"])
    accent = _rgb(_THEME["accent"])

    _add_textbox(slide, 0.9, 2.4, 9.5, 1.2, outline.get("title", "Presentation"),
                 font_size=44, bold=True, color=white)
    _add_textbox(slide, 0.9, 3.7, 8.5, 0.8, outline.get("subtitle", ""),
                 font_size=22, color=accent)
    _add_textbox(slide, 0.9, 6.55, 5.0, 0.4,
                 datetime.now().strftime("%B %Y"), font_size=14, color=_rgb(_THEME["text_muted"]))

    img_path = _try_generate_slide_image(
        outline.get("theme_keywords", "technology AI assistant"), outline.get("title", "")
    )
    if img_path and img_path.exists():
        try:
            slide.shapes.add_picture(str(img_path), Inches(8.8), Inches(1.8), width=Inches(3.8))
        except Exception:
            pass


def _add_content_slide(prs, item: dict) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_decorative_shapes(slide, "content")
    _add_logo(slide, left=12.0, top=0.18, height=0.45)

    white = _rgb(_THEME["white"])
    dark = _rgb(_THEME["text_dark"])
    accent = _rgb(_THEME["accent"])

    title = str(item.get("title", "Slide"))
    _add_textbox(slide, 0.55, 0.22, 11.0, 0.65, title, font_size=28, bold=True, color=white)

    bullets = item.get("bullets") or []
    y = 1.45
    for i, bullet in enumerate(bullets[:6]):
        from pptx.enum.shapes import MSO_SHAPE
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(0.55), Inches(y + 0.08), Inches(0.12), Inches(0.12)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = accent
        diamond.line.fill.background()
        _add_textbox(slide, 0.85, y, 5.8, 0.55, str(bullet), font_size=18, color=dark)
        y += 0.72

    visual = str(item.get("visual_hint") or "")
    img_left = 6.85
    panel_top, panel_w = 1.35, 6.0
    img_path = _try_generate_slide_image(visual, title) if visual and _BUILD_WITH_IMAGES else None
    if img_path and img_path.exists():
        try:
            slide.shapes.add_picture(str(img_path), Inches(img_left), Inches(panel_top), width=Inches(panel_w))
        except Exception:
            _add_visual_png(slide, item, img_left, panel_top, panel_w)
    else:
        _add_visual_png(slide, item, img_left, panel_top, panel_w)


def _add_image_placeholder(slide, left, top, width, height, label: str):
    _draw_visual_panel(slide, {"visual_hint": label, "title": label, "bullets": []}, left, top, width, height)


def _draw_visual_panel(slide, item: dict, left: float, top: float, width: float, height: float) -> None:
    """Draw architecture / stack / icon visuals — never leave an empty placeholder box."""
    title = str(item.get("title") or "").lower()
    hint = str(item.get("visual_hint") or "").lower()
    layout = str(item.get("layout") or "").lower()
    bullets = item.get("bullets") or []

    if layout == "stack" or "stack" in title or "pyramid" in hint or "stack" in hint:
        layers = item.get("stack_layers") or _layers_from_bullets(bullets)
        if layers:
            _draw_stack_pyramid(slide, layers, left, top, width, height)
            return

    if layout == "diagram" or any(k in title for k in ("architecture", "pipeline", "flow")):
        nodes = item.get("diagram_nodes") or _layers_from_bullets(bullets) or ["Input", "Core", "Output"]
        _draw_mini_pipeline(slide, nodes, left, top, width, height)
        return

    if any(k in hint for k in ("security", "shield", "lock")):
        _draw_icon_badge(slide, "Security", left, top, width, height)
    elif any(k in hint for k in ("devops", "ci/cd", "pipeline")):
        _draw_devops_icons(slide, left, top, width, height)
    elif any(k in hint for k in ("cloud", "kubernetes", "infra")):
        _draw_cloud_stack(slide, left, top, width, height)
    else:
        _draw_icon_grid(slide, bullets or ["AI", "Voice", "Tools"], left, top, width, height)


def _draw_stack_pyramid(slide, layers: list[str], left, top, width, height) -> None:
    """Technology stack pyramid — stacked layers from bullets."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    n = max(len(layers), 1)
    _add_rect(slide, left, top, width, height, _rgb(_THEME["white"]), _rgb(_THEME["accent"]))
    gap = 0.08
    layer_h = (height - gap * (n + 1)) / n
    white = _rgb(_THEME["white"])

    for i, label in enumerate(layers[:6]):
        shrink = i * 0.11
        lw = max(width - shrink * 2, width * 0.45)
        lx = left + (width - lw) / 2
        ly = top + gap + i * (layer_h + gap)
        colors = [_rgb(_THEME["accent"]), _rgb(_THEME["accent2"]), _rgb(_THEME["bg_dark"])]
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ly), Inches(lw), Inches(layer_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.color.rgb = white
        tf = box.text_frame
        tf.paragraphs[0].text = label[:36]
        tf.paragraphs[0].font.size = Pt(12 if n > 4 else 14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = white


def _draw_mini_pipeline(slide, nodes: list[str], left, top, width, height) -> None:
    """Vertical or horizontal pipeline inside content panel."""
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.util import Inches, Pt

    _add_rect(slide, left, top, width, height, _rgb((0x12, 0x2B, 0x45)), _rgb(_THEME["accent"]))
    n = min(len(nodes), 5)
    if n <= 0:
        return
    box_h = min(0.55, (height - 0.4) / n - 0.12)
    box_w = width - 0.6
    x = left + 0.3
    white = _rgb(_THEME["white"])
    accent = _rgb(_THEME["accent"])
    cy_list: list[float] = []

    for i, label in enumerate(nodes[:n]):
        y = top + 0.25 + i * (box_h + 0.18)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(box_h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = accent if i % 2 == 0 else _rgb(_THEME["accent2"])
        box.line.color.rgb = white
        tf = box.text_frame
        tf.paragraphs[0].text = label[:32]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = white
        cy_list.append(y + box_h / 2)

    for i in range(len(cy_list) - 1):
        y1, y2 = cy_list[i], cy_list[i + 1]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x + box_w / 2), Inches(y1 + box_h / 2),
            Inches(x + box_w / 2), Inches(y2 - box_h / 2),
        )
        conn.line.color.rgb = accent
        conn.line.width = Pt(2)


def _draw_icon_grid(slide, labels: list, left, top, width, height) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    _add_rect(slide, left, top, width, height, _rgb(_THEME["bg_light"]), _rgb(_THEME["accent"]))
    items = [str(x).split(":")[0][:18] for x in labels[:4]] or ["AI", "Data", "Cloud", "Ops"]
    cols = 2
    cw, ch = width / cols - 0.15, height / 2 - 0.2
    for i, label in enumerate(items):
        col, row = i % cols, i // cols
        cx = left + 0.1 + col * (cw + 0.1)
        cy = top + 0.15 + row * (ch + 0.15)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = _rgb(_THEME["accent"])
        circ.line.fill.background()
        _add_textbox(slide, cx - 0.05, cy + 0.62, cw + 0.1, 0.45, label, font_size=11,
                     bold=True, color=_rgb(_THEME["text_dark"]))


def _draw_icon_badge(slide, label: str, left, top, width, height) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    _add_rect(slide, left, top, width, height, _rgb(_THEME["bg_dark"]), _rgb(_THEME["accent"]))
    shield = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON, Inches(left + width / 2 - 0.7), Inches(top + 0.8), Inches(1.4), Inches(1.4)
    )
    shield.fill.solid()
    shield.fill.fore_color.rgb = _rgb(_THEME["accent"])
    shield.line.color.rgb = _rgb(_THEME["white"])
    _add_textbox(slide, left + 0.3, top + height - 1.2, width - 0.6, 0.8, label,
                 font_size=16, bold=True, color=_rgb(_THEME["white"]))


def _draw_devops_icons(slide, left, top, width, height) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    _add_rect(slide, left, top, width, height, _rgb((0x12, 0x2B, 0x45)), _rgb(_THEME["accent"]))
    steps = ["Build", "Test", "Deploy", "Monitor"]
    sw = (width - 0.5) / len(steps)
    for i, step in enumerate(steps):
        sx = left + 0.25 + i * sw
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(sx), Inches(top + height / 2 - 0.35), Inches(sw - 0.08), Inches(0.7)
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = _rgb(_THEME["accent"]) if i % 2 == 0 else _rgb(_THEME["accent2"])
        rect.line.fill.background()
        tf = rect.text_frame
        tf.paragraphs[0].text = step
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = _rgb(_THEME["white"])


def _draw_cloud_stack(slide, left, top, width, height) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    _add_rect(slide, left, top, width, height, _rgb(_THEME["bg_light"]), _rgb(_THEME["accent"]))
    for i, label in enumerate(["K8s", "Cloud", "Scale"]):
        cy = top + 0.5 + i * 1.5
        cloud = slide.shapes.add_shape(
            MSO_SHAPE.CLOUD, Inches(left + 1.5), Inches(cy), Inches(2.5), Inches(1.0)
        )
        cloud.fill.solid()
        cloud.fill.fore_color.rgb = _rgb(_THEME["accent2"])
        cloud.line.color.rgb = _rgb(_THEME["accent"])
        _add_textbox(slide, left + 4.2, cy + 0.25, 1.5, 0.5, label, font_size=14, bold=True,
                     color=_rgb(_THEME["text_dark"]))


def _add_stack_slide(prs, item: dict) -> None:
    """Full-slide technology stack with pyramid + bullet details."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_decorative_shapes(slide, "content")
    _add_logo(slide, left=12.0, top=0.18, height=0.45)

    white = _rgb(_THEME["white"])
    dark = _rgb(_THEME["text_dark"])
    accent = _rgb(_THEME["accent"])

    title = str(item.get("title", "Stack"))
    _add_textbox(slide, 0.55, 0.22, 11.0, 0.65, title, font_size=28, bold=True, color=white)

    layers = item.get("stack_layers") or _layers_from_bullets(item.get("bullets") or [])
    stack_item = {**item, "stack_layers": layers or ["Layer 1", "Layer 2", "Layer 3"]}
    _add_visual_png(slide, stack_item, 6.85, 1.35, 6.0)

    bullets = item.get("bullets") or []
    y = 1.45
    for bullet in bullets[:5]:
        from pptx.enum.shapes import MSO_SHAPE
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(0.55), Inches(y + 0.08), Inches(0.12), Inches(0.12)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = accent
        diamond.line.fill.background()
        _add_textbox(slide, 0.85, y, 5.8, 0.55, str(bullet), font_size=16, color=dark)
        y += 0.78


def _add_diagram_slide(prs, item: dict) -> None:
    from actions.presentation_graphics import render_architecture_flow

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_decorative_shapes(slide, "content")
    _add_logo(slide, left=12.0, top=0.18, height=0.45)

    title = str(item.get("title", "Diagram"))
    _add_textbox(slide, 0.55, 0.22, 11.0, 0.65, title, font_size=28, bold=True, color=_rgb(_THEME["white"]))

    nodes = [str(n) for n in (item.get("diagram_nodes") or item.get("bullets") or [])[:6]]
    if len(nodes) < 3:
        nodes = ["Voice / UI", "JARVIS Core", "Tools & APIs", "Response"]

    try:
        png = render_architecture_flow(nodes, title="")
        _embed_png(slide, png, 0.45, 1.25, width=12.4)
    except Exception as e:
        print(f"[PresentationCreate] Diagram PNG failed: {e}")
        _draw_mini_pipeline(slide, nodes, 1.0, 1.4, 11.5, 5.2)

    bullets = item.get("bullets") or []
    if bullets:
        _add_textbox(slide, 0.55, 6.72, 12.0, 0.45, str(bullets[0])[:120],
                     font_size=13, color=_rgb(_THEME["text_muted"]))


def _add_section_slide(prs, item: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, 13.333, 7.5, _rgb(_THEME["accent"]))
    _add_logo(slide, left=11.8, top=0.25, height=0.55)
    _add_textbox(
        slide, 1.0, 2.8, 11.0, 1.5, str(item.get("title", "Section")),
        font_size=40, bold=True, color=_rgb(_THEME["white"]),
    )


def _add_closing_slide(prs, item: dict, outline: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_decorative_shapes(slide, "title")
    _add_logo(slide, left=5.9, top=1.2, height=1.0)

    white = _rgb(_THEME["white"])
    accent = _rgb(_THEME["accent"])
    _add_textbox(slide, 1.0, 3.2, 11.3, 0.9, "Thank You", font_size=40, bold=True, color=white)

    bullets = item.get("bullets") or []
    if bullets:
        summary = "  •  ".join(str(b) for b in bullets[:4])
        _add_textbox(slide, 1.0, 4.3, 11.3, 1.0, summary, font_size=16, color=accent)

    _add_textbox(slide, 1.0, 6.3, 11.0, 0.4, outline.get("subtitle", ""), font_size=14, color=_rgb(_THEME["text_muted"]))


def _build_pptx(outline: dict, output_path: Path, *, with_images: bool = True) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    global _BUILD_WITH_IMAGES
    _BUILD_WITH_IMAGES = with_images

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, outline)

    for item in outline["slides"]:
        layout = str(item.get("layout") or "content").lower()
        if layout == "diagram":
            _add_diagram_slide(prs, item)
        elif layout == "stack":
            _add_stack_slide(prs, item)
        elif layout == "section":
            _add_section_slide(prs, item)
        elif layout == "closing":
            _add_closing_slide(prs, item, outline)
        else:
            _add_content_slide(prs, item)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _open_path(path: Path) -> str:
    if platform.system() == "Darwin":
        subprocess.run(["open", str(path)], check=False, timeout=15)
        return f"Opened {path.name}."
    if platform.system() == "Windows":
        subprocess.run(["start", "", str(path)], shell=True, check=False)
        return f"Opened {path.name}."
    subprocess.run(["xdg-open", str(path)], check=False)
    return f"Opened {path.name}."


def _create_local_presentation(
    title: str,
    topic: str,
    num_slides: int,
    audience: str,
    output_path: Path,
    *,
    open_after: bool = True,
    with_images: bool = True,
) -> str:
    outline = _generate_outline(title, topic, num_slides, audience)
    _build_pptx(outline, output_path, with_images=with_images)

    slide_titles = [str(s.get("title", "")) for s in outline.get("slides", [])]
    resolved = output_path.resolve()
    summary = (
        f"[PRESENTATION: ok]\n"
        f"Created professional presentation '{outline.get('title', title)}' "
        f"with {len(slide_titles)} slides.\n"
        f"Path: {resolved}\n"
        f"Slides: {', '.join(slide_titles[:8])}"
        + ("…" if len(slide_titles) > 8 else "")
        + "\n\nStyled with JARVIS corporate theme (logo, typography, Pillow architecture diagrams)."
        + "\nUpload to Google Drive: drag the .pptx to Drive, then Open with → Google Slides."
    )

    if open_after:
        try:
            summary += f"\n{_open_path(output_path)}"
        except Exception as e:
            summary += f"\nCould not open file: {e}"

    return summary


def presentation_create(
    parameters: dict | None = None,
    player=None,
) -> str:
    """
    Actions:
      create          — professional local .pptx (default)
      create_in_drive — same as create (local only; upload to Drive manually)
      open_slides     — instructions only (no browser automation)
    """
    params = parameters or {}
    action = (params.get("action") or "create").lower().strip()

    title = (params.get("title") or "Presentation").strip()
    topic = (
        params.get("topic")
        or params.get("outline")
        or params.get("description")
        or params.get("brief")
        or ""
    ).strip()
    num_slides = min(max(int(params.get("num_slides") or 8), 3), 20)
    audience = (params.get("audience") or "").strip()
    open_after = bool(params.get("open_after", True))
    with_images = bool(params.get("with_images", False))

    if action == "open_slides":
        return (
            "Local presentation mode only. Ask JARVIS to create a presentation — "
            "then upload the .pptx from ~/Documents/JARVIS Presentations/ to Google Drive."
        )

    # Roll back Drive automation — always create locally
    if action in ("create_in_drive", "google_drive", "drive", "blank"):
        if not topic and action != "blank":
            return "Provide a topic for the presentation."
        if action == "blank":
            return (
                "[PRESENTATION: info]\n"
                "Blank Google Slides is not automated. "
                "Say 'create a presentation about …' for a professional local .pptx."
            )
        action = "create"

    if action != "create":
        return f"Unknown action '{action}'. Use: create"

    if not topic:
        return "Provide a topic, outline, or description for the presentation."

    out_raw = (params.get("output_path") or "").strip()
    if out_raw:
        output_path = Path(out_raw).expanduser()
        if output_path.is_dir():
            output_path = output_path / _safe_filename(title)
        elif not str(output_path).lower().endswith(".pptx"):
            output_path = output_path.with_suffix(".pptx")
    else:
        _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
        stamp = datetime.now().strftime("%Y%m%d_%H%M")
        output_path = _OUTPUT_DIR / f"{stamp}_{_safe_filename(title)}"

    log = f"[PresentationCreate] {title!r} → {output_path.name} ({num_slides} slides, professional)"
    print(log)
    if player and hasattr(player, "write_log"):
        player.write_log(log)
        player.write_log("[presentation] building slides — please wait…")

    try:
        msg = _create_local_presentation(
            title, topic, num_slides, audience, output_path,
            open_after=open_after, with_images=with_images,
        )
    except Exception as e:
        return f"Presentation creation failed: {e}"

    if player and hasattr(player, "show_content"):
        player.show_content(f"PRESENTATION — {title}", msg[:800])
    if player and hasattr(player, "write_log"):
        player.write_log("[presentation] local professional .pptx created")

    return msg
